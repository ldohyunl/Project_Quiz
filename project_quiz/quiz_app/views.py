from django.shortcuts import render
from .forms import FileUploadForm
import os
import openai
import fitz  # PyMuPDF (PDF용)
import openpyxl  # Excel용
from pptx import Presentation  # PPT용
from docx import Document  # Word(DOCX)용

# OpenAI API 클라이언트 초기화
client = openai.OpenAI(api_key="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")

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = [shape.text.strip() for slide in presentation.slides for shape in slide.shapes if
            hasattr(shape, "text") and shape.text.strip()]
    return "\n".join(text)

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = [page.get_text("text") for page in doc]
    return "\n".join(text)

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    return "\n".join(text)

def extract_text_from_xlsx(file_path):
    workbook = openpyxl.load_workbook(file_path)
    text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            text.append(" | ".join([str(cell.value) for cell in row if cell.value]))
    return "\n".join(text)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    else:
        raise ValueError("지원되지 않는 파일 형식입니다.")

def generate_quiz_with_gpt(text):
    prompt = f"""
    Ask 10 questions in Korean based on the provided text format.

    === 파일 내용 ===
    {text}

    🔹 문제 형식:
    1. 질문 내용?
       a) 선택지 1
       b) 선택지 2
       c) 선택지 3
       d) 선택지 4
       정답: [정답 선택지]
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You have excellent knowledge in all fields. Create difficult problems in a format based on text provided by the user."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=3000
    )
    return response.choices[0].message.content

def index(request):
    form = FileUploadForm()
    if request.method == "POST":
        form = FileUploadForm(request.POST, request.FILES)
        if form.is_valid():
            file = form.cleaned_data["file"]
            # 업로드 디렉토리 경로 설정
            upload_folder = "uploads"
            if not os.path.exists(upload_folder):
                os.makedirs(upload_folder)

            file_path = os.path.join(upload_folder, file.name)
            with open(file_path, "wb") as f:
                for chunk in file.chunks():
                    f.write(chunk)

            try:
                # 텍스트 추출 후 퀴즈 생성 요청
                file_text = extract_text(file_path)
                if file_text:
                    quiz = generate_quiz_with_gpt(file_text)
                    return render(request, 'quiz_app/index.html', {'quiz': quiz, 'form': form})
                else:
                    return render(request, 'quiz_app/index.html', {'error': "파일에서 텍스트를 추출할 수 없습니다.", 'form': form})
            except Exception as e:
                return render(request, 'quiz_app/index.html', {'error': f"오류 발생: {e}", 'form': form})
    return render(request, 'quiz_app/index.html', {'form': form})